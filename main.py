# Adding more details to each slide in the presentation
from pptx import Presentation
# Create a new presentation object with more detailed content
prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]  # Title slide layout
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Introduction to C/C++"
subtitle.text = "Vivek Sharma\nOctober 2024"

# Slide 1: Definitions
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "Definitions"

content = slide.shapes.placeholders[1].text_frame
content.text = "C:\nA procedural programming language that provides low-level memory access. Known for its efficiency and close-to-hardware features.\n\nC++:\nAn extension of C that incorporates object-oriented features, allowing for a broader range of applications, including real-world modeling."

# Slide 2: History
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "History"

content = slide.shapes.placeholders[1].text_frame
content.text = "• C was developed by Dennis Ritchie at Bell Labs in 1972 to rewrite Unix.\n• C++ was developed by Bjarne Stroustrup in 1985, building upon C by adding object-oriented programming features."

# Slide 3: Key Features of C
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "Key Features of C"

content = slide.shapes.placeholders[1].text_frame
content.text = "• Procedural Programming: Functions drive the program's flow.\n• Low-level Access: Direct access to memory using pointers.\n• Portability: C programs can run on many different machines.\n• Efficient: C provides fine control over system resources and hardware."

# Slide 4: Key Features of C++
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "Key Features of C++"

content = slide.shapes.placeholders[1].text_frame
content.text = "• Object-Oriented Programming (OOP): C++ supports classes and objects to represent real-world entities.\n• Standard Template Library (STL): A powerful set of C++ libraries for data structures and algorithms.\n• Operator Overloading: Allows customizing how operators work for user-defined types.\n• Inheritance and Polymorphism: Key OOP features to enhance code reuse and flexibility."

# Slide 5: Syntax Comparison
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "Syntax Comparison"

content = slide.shapes.placeholders[1].text_frame
content.text = "• C uses functions, while C++ allows defining classes and objects.\n• Memory management in C is done via malloc()/free(), while C++ uses new/delete.\n• C++ supports function overloading, while C does not."

# Slide 6: Memory Management
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "Memory Management"

content = slide.shapes.placeholders[1].text_frame
content.text = "• C uses malloc() and free() for dynamic memory management.\n• C++ introduces new and delete operators for allocating and deallocating memory.\n• Stack memory vs Heap memory: Local variables use stack memory, while dynamically allocated variables use heap memory."

# Slide 7: Common Use Cases
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "Common Use Cases"

content = slide.shapes.placeholders[1].text_frame
content.text = "• C is widely used in system programming, embedded systems, and low-level hardware programming.\n• C++ is preferred for game development, graphical user interfaces (GUIs), real-time systems, and applications that require OOP."

# Slide 8: Advantages and Disadvantages
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "Advantages and Disadvantages"

content = slide.shapes.placeholders[1].text_frame
content.text = "Advantages of C:\n• Low-level programming capabilities\n• Portability\n\nDisadvantages of C:\n• No object-oriented features\n• Limited error handling\n\nAdvantages of C++:\n• OOP supports modular code and reusability\n• STL and other libraries make development faster\n\nDisadvantages of C++:\n• Complex syntax compared to C\n• Slower compilation times."

# Slide 9: Development Tools
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "Development Tools"

content = slide.shapes.placeholders[1].text_frame
content.text = "• Code::Blocks: A free, open-source IDE supporting C/C++ development.\n• Visual Studio: A comprehensive IDE from Microsoft with debugging and testing features.\n• GCC: The GNU Compiler Collection, commonly used for compiling C and C++.\n• Clang: A modern C/C++ compiler focusing on performance and diagnostics."

# Slide 10: Best Practices
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "Best Practices"

content = slide.shapes.placeholders[1].text_frame
content.text = "• Follow standard naming conventions for variables and functions.\n• Write modular code for maintainability and reusability.\n• Use comments and documentation to improve code readability.\n• Avoid memory leaks by ensuring all dynamically allocated memory is freed."

# Slide 11: Resources for Further Learning
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "Resources for Further Learning"

content = slide.shapes.placeholders[1].text_frame
content.text = "Books:\n• 'The C Programming Language' by Brian W. Kernighan and Dennis M. Ritchie\n• 'C++ Primer' by Stanley Lippman\n\nWebsites:\n• cppreference.com: A comprehensive online reference for C++.\n• cplusplus.com: A website providing C++ tutorials, references, and resources."

# Slide 12: Q&A
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "Q&A"

content = slide.shapes.placeholders[1].text_frame
content.text = "Feel free to ask any questions."

# Save the presentation with more detailed content
detailed_ppt_path = "/Introduction_to_C_and_CPP_Detailed_Presentation.pptx"
prs.save(detailed_ppt_path)

detailed_ppt_path
